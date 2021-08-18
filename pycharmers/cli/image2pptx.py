# coding: utf-8
import argparse
import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from tqdm import tqdm

IMAGE_EXTENSIONS = [".jpg", ".png", ".jpeg"]

def image2pptx(argv=sys.argv[1:]):
    """Paste images to PowerPoint.

    Args:
        --image-path (Path, optional)     : Paths to image files. Defaults to ``()``.
        --image-dir (Path, optional)      : Path to the directory where images are. Defaults to ``None``.
        -W/--slide-width (int, optional)  : The width of PowerPoint slide. Defaults to ``9144000``.
        -H/--slide-height (int, optional) : The height of PowerPoint slide. Defaults to ``6858000``.
        --slide-size (str, optional)      : The size of PowerPoint slide. Please chose from ``["4:3", "16:9"]``. Defaults to ``"4:3"``.
        -O/--outpptx (Path, optional)     : The path to the created PowerPoint. Defaults to ``Path("test.pptx")``.

    Note:
        When you run from the command line, execute as follows::

            $ image2pptx --image-path /path/to/image1.png \\
                                      /path/to/image2.jpg \\
                                      /path/to/image3.jpeg \\
                         --image_dir /path/to/image_dir \\
                         --slide-size "16:9" \\
                         --outpptx "image.pptx"
    """
    parser = argparse.ArgumentParser(prog="image2pptx", description="Paste images to PowerPoint", add_help=True)
    parser.add_argument("--image-path", type=Path, nargs="*", help="Paths to image files.")
    parser.add_argument("--image-dir", type=Path, help="Path to the directory where images are.")
    parser.add_argument("-W", "--slide-width", type=int, default=9144000, help="The width of PowerPoint slide.")
    parser.add_argument("-H", "--slide-height", type=int, default=6858000, help="The height of PowerPoint slide.")
    parser.add_argument("--slide-size", type=str, default=None, choices=["4:3", "16:9"], help="The size of PowerPoint slide.")
    parser.add_argument("-O", "--outpptx", type=Path, default=Path("test.pptx"), help="The path to the created PowerPoint.")
    args = parser.parse_args()

    image_paths = list(args.image_path)
    if args.image_dir is not None:
        image_paths += sorted(
            [
                path
                for path in args.image_dir.rglob("*")
                if path.suffix.lower() in IMAGE_EXTENSIONS
            ]
        )

    slide_size = args.slide_size
    if slide_size is not None:
        slide_width, slide_height = {
            "4:3": (9144000, 6858000),
            "16:9": (12193200, 6858000),
        }[slide_size]
    else:
        slide_width = args.slide_width
        slide_height = args.slide_height

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = 0
    for image_path in image_paths:
        img = Image.open(image_path)
        slide.shapes.add_picture(
            image_file=str(image_path),
            left=left,
            top=top,
        )
        left += img.width * 1e4
        if left >= slide_width:
            top += img.height * 1e4
            left = 0

    prs.save(file=args.outpptx)
